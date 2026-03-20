import streamlit as st
import time
import os
import requests
import base64
from PIL import Image
import torch
from transformers import CLIPProcessor, CLIPModel
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.common.exceptions import InvalidArgumentException
import platform
import numpy as np
from sklearn.cluster import KMeans
from pptx import Presentation
from pptx.util import Inches, Pt
import io

@st.cache_resource
def load_model():
    print("AI 엔진(CLIP) 다운로드 및 로딩을 시작합니다... (최초 1회 소요)")
    processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
    model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
    return model, processor

def create_pptx_report(keyword, category_dynamic_data, classified_images):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"AI Product Design Analysis: {keyword}"
    subtitle.text = "Pinterest Reference Scraping & CMF Analysis Report"
    
    bullet_slide_layout = prs.slide_layouts[5] # 빈 슬라이드 + 제목
    
    for category, data in category_dynamic_data.items():
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = f"{category} Style Analysis"
        
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(1)
        txBox = shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        features_str = ", ".join(data.get("features", []))
        p = tf.add_paragraph()
        p.text = f"✨ Core Features (CMF): {features_str}"
        p.font.size = Pt(14)
        
        colors_str = ", ".join(data.get("hex_colors", []))
        if colors_str:
            p2 = tf.add_paragraph()
            p2.text = f"🎨 Extracted Color Palette: {colors_str}"
            p2.font.size = Pt(14)
            
        prompt = data.get("prompt", "")
        if prompt:
            p3 = tf.add_paragraph()
            p3.text = f"✍️ Suggested AI Prompt:\n{prompt}"
            p3.font.size = Pt(11)
            
        img_list = classified_images.get(category, [])
        img_left = Inches(0.5)
        img_top = Inches(3.8) # 텍스트 영역 아래부터
        img_width = Inches(1.7)
        row_height = Inches(1.8)
        col_width = Inches(1.8)
        
        items_on_current_slide = 0
        max_items_first_slide = 10 # 타이틀 슬라이드는 하단 2줄
        max_items_cont_slide = 15 # 연장 슬라이드부터는 상단부터 3줄
        
        is_first_slide = True
        
        for img_path in img_list: # [:4] 개수 제한 해제 (전체 복사)
            if is_first_slide and items_on_current_slide >= max_items_first_slide:
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                shapes.title.text = f"{category} Style Analysis (Continued)"
                img_left = Inches(0.5)
                img_top = Inches(1.5)
                items_on_current_slide = 0
                is_first_slide = False
            elif not is_first_slide and items_on_current_slide >= max_items_cont_slide:
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                shapes.title.text = f"{category} Style Analysis (Continued)"
                img_left = Inches(0.5)
                img_top = Inches(1.5)
                items_on_current_slide = 0
            
            try:
                shapes.add_picture(img_path, img_left, img_top, width=img_width)
            except:
                pass
            
            items_on_current_slide += 1
            img_left += col_width
            
            # 한 줄에 5개가 차면 줄바꿈
            if items_on_current_slide % 5 == 0: 
                img_left = Inches(0.5)
                img_top += row_height

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# 앱 기본 설정
st.set_page_config(page_title="RefMaster AI - 제품 디자인 탐색기", layout="wide")

# 매끈하고 모던한 UI 디자인 적용 (기능 영향 절대 없음)
st.markdown("""
<style>
    .hero-title {
        font-size: 3.5rem !important;
        font-weight: 900 !important;
        letter-spacing: -1.5px;
        margin-bottom: 0.5rem !important;
        line-height: 1.2;
        background: linear-gradient(135deg, #1A1A1A 0%, #FF4B2B 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    .hero-subtitle {
        font-size: 1.3rem;
        color: #555;
        font-weight: 500;
        margin-bottom: 2rem;
        word-break: keep-all;
        line-height: 1.5;
    }
    .panel-box {
        background-color: #f8f9fa;
        padding: 24px;
        border-radius: 16px;
        border: 1px solid #eaeaea;
        margin-bottom: 1rem;
    }
    /* 시작 버튼 디자인 강화 */
    div.stButton > button {
        height: 65px;
        font-size: 22px !important;
        font-weight: 800 !important;
        border-radius: 12px;
        border: none;
        transition: all 0.3s ease;
    }
    div.stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 20px rgba(255, 75, 43, 0.2);
    }
</style>
<h1 class="hero-title">🏆 AI 제품 디자인 분류 마스터</h1>
<p class="hero-subtitle">키워드 하나만 입력하면 핀터레스트에서 완벽한 레퍼런스를 수집하고, <br>조형적 CMF(형태/재질/색상) 데이터를 딥러닝으로 자동 분석해 드립니다.</p>
""", unsafe_allow_html=True)

# 모델 로드 (앱 상단에 고정 표시)
model_load_state = st.empty()
with model_load_state.container():
    with st.spinner("🧠 딥러닝 기반 제품 디자인 조형 분석 엔진을 깨우는 중..."):
        model, processor = load_model()
    st.success("✅ AI 디자인 분석 엔진 준비 완료!")
    time.sleep(1)
model_load_state.empty()

# ==========================================
# 사용자 설정 영역
# ==========================================
st.markdown("<h2 style='font-weight:800; font-size:1.8rem; margin-top:10px; margin-bottom:15px;'>🔍 레퍼런스 탐색 및 디자인 분석 설정</h2>", unsafe_allow_html=True)

col1, col2 = st.columns([1.2, 1])
with col1:
    keyword = st.text_input("💡 수집할 제품 키워드 (영어 검색 권장)", "Bluetooth speaker design")
    target_count = st.number_input("🎯 수집할 레퍼런스 누적 목표 개수", min_value=1, max_value=1000, value=50, step=10)
with col2:
    st.write("✨ **AI 조형 자동 분류 (Zero-Shot)**")
    st.info("AI 모델이 25가지 특화된 디자인 언어 및 CMF 사전을 기반으로 레퍼런스의 핵심 스타일을 스스로 분류합니다.")
categories = [
    "Minimalist", "Industrial", "Vintage", "Cyberpunk", "Organic modern", 
    "Retro futuristic", "Scandinavian", "Braun style", "Apple minimalist", 
    "Steampunk", "Mid-Century Modern", "Bauhaus", "Brutalist", 
    "High-tech", "Art Deco", "Futuristic", "Biophilic", "Space Age",
    "Y2K aesthetic", "Postmodern", "Pop Art", "Wabi-sabi", 
    "Surrealist", "Tactical", "Luxurious"
]

st.write("---")

# ==========================================
# 세션 상태 초기화
# ==========================================
if "analysis_complete" not in st.session_state:
    st.session_state.analysis_complete = False
if "saved_files" not in st.session_state:
    st.session_state.saved_files = []
if "classified_images" not in st.session_state:
    st.session_state.classified_images = {}
if "used_keyword" not in st.session_state:
    st.session_state.used_keyword = ""

# ==========================================
# 실행 영역
# ==========================================
if st.button("🚀 AI 레퍼런스 탐색 및 딥러닝 분석 시작하기", type="primary", use_container_width=True):
    if len(categories) == 0:
        st.error("❌ 분류 기준을 1개 이상 입력해 주세요.")
    else:
        st.session_state.analysis_complete = False
        st.info(f"📍 총 {len(categories)}개의 분류 기준으로 AI 분석을 시작합니다: {', '.join(categories)}")
        # 저장할 폴더 자동 생성 (현재 폴더 아래)
        save_path = os.path.join(os.getcwd(), f"References_{keyword.replace(' ', '_')}")
        if not os.path.exists(save_path):
            os.makedirs(save_path)
            
        st.info("💡 단계를 시작합니다! 핀터레스트 로봇이 작동하는 동안 열린 크롬 창을 스스로 닫지 마세요.")
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        # 1단계: 핀터레스트 크롤링 및 다운로드
        status_text.write("🔍 [1/2 단계] 인터넷에서 고화질 레퍼런스 고속 수집 중...")
        
        driver = None
        try:
            options = webdriver.ChromeOptions()
            is_linux = platform.system() == "Linux"
            
            if is_linux:
                # 1. 서버(Linux) 환경용 설정
                options.add_argument("--headless=new") 
                options.add_argument("--no-sandbox")
                options.add_argument("--disable-dev-shm-usage")
                options.add_argument("--disable-gpu")
                options.binary_location = "/usr/bin/chromium"
                try:
                    driver = webdriver.Chrome(service=Service("/usr/bin/chromedriver"), options=options)
                except Exception:
                    driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=options)
            else:
                # 2. 로컬(Windows) 환경용 설정
                # 디버깅을 위해 로컬에서는 화면을 띄움 (원하시면 --headless 추가 가능)
                # options.add_argument("--headless") 
                bot_data_dir = os.path.join(os.getcwd(), 'bot_chrome_profile')
                options.add_argument(f"user-data-dir={bot_data_dir}")
                options.add_argument("--disable-session-crashed-bubble")
                options.add_experimental_option("excludeSwitches", ["enable-automation"])
                driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()), options=options)
        except Exception as e:
            st.error(f"🚨 크롬 로드 실패: {e}")
            st.stop()
                
        try:
            # 브라우저 열고 이동
            time.sleep(1)
            driver.get(f"https://www.pinterest.com/search/pins/?q={keyword}")
            status_text.write("⏳ [알림] 봇이 핀터레스트에 접속했습니다. 15초 대기 중...")
            time.sleep(15) 
            
            seen_srcs = set()
            saved_files = [] 
            no_new_image_attempts = 0 
            
            with st.spinner("최고의 레퍼런스를 핀터레스트에서 추출 중입니다..."):
                while len(saved_files) < target_count:
                    driver.execute_script("window.scrollTo(0, document.body.scrollHeight);")
                    time.sleep(3) 
                    
                    try:
                        imgs = driver.find_elements(By.TAG_NAME, "img")
                    except Exception:
                        st.error("⚠️ 로봇 작동 중 크롬 창이 의도치 않게 닫혔습니다!")
                        break
                        
                    new_found_this_turn = 0
                    for img in reversed(imgs): 
                        if len(saved_files) >= target_count: break
                        try:
                            width = img.size['width']
                            height = img.size['height']
                            src = img.get_attribute("src")
                            
                            if width > 120 and height > 120 and src and "pinimg.com" in src and "profile" not in src and src not in seen_srcs:
                                high_res_src = src.replace("/236x/", "/736x/")
                                try:
                                    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"}
                                    img_data = requests.get(high_res_src, headers=headers).content
                                    file_name = f"ref_{len(saved_files)+1}.png"
                                    filepath = os.path.join(save_path, file_name)
                                    with open(filepath, 'wb') as handler:
                                        handler.write(img_data)
                                    seen_srcs.add(src)
                                    saved_files.append(filepath)
                                    new_found_this_turn += 1
                                    progress_bar.progress((len(saved_files) / target_count) * 0.5)
                                    status_text.write(f"📥 [1/2 단계] {len(saved_files)}/{target_count} 장 사진 다운로드 성공!")
                                except Exception: continue
                        except Exception: continue
                    
                    if new_found_this_turn == 0:
                        no_new_image_attempts += 1
                        if no_new_image_attempts >= 6: break
                    else:
                        no_new_image_attempts = 0
                            
            if driver:
                driver.quit()
            
            if len(saved_files) == 0:
                st.warning("수집된 이미지가 없습니다. 다시 시도해보세요.")
                st.stop()
                
            # 2단계: AI 분석
            status_text.write(f"🧠 [2/2 단계] 수집된 {len(saved_files)}장의 레퍼런스 분석 시작...")
            classified_images = {c: [] for c in categories}
            
            with st.spinner("AI 분석 중..."):
                for idx, filepath in enumerate(saved_files):
                    try:
                        pil_img = Image.open(filepath).convert("RGB")
                        prompt_templates = [f"A professional product photo of a {keyword} in {category} style" for category in categories]
                        inputs = processor(text=prompt_templates, images=pil_img, return_tensors="pt", padding=True)
                        with torch.no_grad():
                            outputs = model(**inputs)
                        probs = outputs.logits_per_image.softmax(dim=1)
                        best_category = categories[probs.argmax().item()]
                        classified_images[best_category].append(filepath)
                        progress_bar.progress(0.5 + ((idx + 1) / len(saved_files)) * 0.5)
                        status_text.write(f"🔄 [2/2 단계] {idx + 1}/{len(saved_files)} 분석 완료: **{best_category}**")
                    except Exception: pass
            
            st.success(f"🎉 환상적입니다! {len(saved_files)}장의 분석을 모두 마쳤습니다.")
            st.balloons()
            
            # 분석 결과를 세션에 저장하여 다운로드나 페이지 리렌더링 시 보존
            st.session_state.saved_files = saved_files
            st.session_state.classified_images = classified_images
            st.session_state.used_keyword = keyword
            st.session_state.analysis_complete = True
            
        except Exception as e:
            if driver: driver.quit()
            st.error(f"실행 중 오류 발생: {e}")

# ==========================================
# 결과 출력 및 다운로드 영역 (st.button 바깥으로 분리)
# ==========================================
if st.session_state.analysis_complete:
    saved_files = st.session_state.saved_files
    classified_images = st.session_state.classified_images
    keyword = st.session_state.used_keyword
    
    # 기존 들여쓰기를 유지하기 위한 구조
    try:
        if True:
            st.header("✨ AI 디자인 스타일 무드보드")
            # 이미지가 할당된 카테고리만 필터링하고 이미지 개수 기준 내림차순 정렬
            active_categories = [c for c in categories if len(classified_images[c]) > 0]
            active_categories.sort(key=lambda c: len(classified_images[c]), reverse=True)
            
            # 수집된 사진에서 가장 대표적인 분류기준을 최대 5개(3~5개)로 제한
            if len(active_categories) > 5:
                active_categories = active_categories[:5]
            
            category_dynamic_data = {}
            feature_pool = [
                "metallic surface", "wooden texture", "glass material", "matte finish", "glossy reflection",
                "curved ergonomic shape", "sharp angular edges", "strict geometric form", "organic flowing lines",
                "translucent parts", "heavy blocky structure", "sleek and slim profile",
                "rough textured surface", "perfectly smooth surface", "leather details", "fabric mesh",
                "soft diffused lighting", "harsh dramatic shadows", "highly colorful", "monochrome palette",
                "complex mechanical details", "clean unadorned surfaces", "rugged tactical build", "elegant premium feel"
            ]
            
            for category in active_categories:
                img_list = classified_images[category]
                st.subheader(f"🎨 {category} 스타일 ({len(img_list)}장)")
                cols = st.columns(4)
                for i, img_path in enumerate(img_list):
                    with cols[i % 4]:
                        with open(img_path, "rb") as img_file:
                            b64 = base64.b64encode(img_file.read()).decode()
                        st.markdown(f'''<a href="data:image/png;base64,{b64}" target="_blank"><img src="data:image/png;base64,{b64}" style="width:100%; border-radius:12px;"></a>''', unsafe_allow_html=True)
                
                # ----------------
                # 수집된 사진들에서 실제 조형적 특징(Feature) 동적 추출 (CLIP Zero-shot)
                # ----------------
                dynamic_features = []
                try:
                    sample_imgs = []
                    for filepath in img_list[:5]: # 너무 오래 걸리지 않게 최대 5장만 샘플링
                        try:
                            sample_imgs.append(Image.open(filepath).convert("RGB"))
                        except: pass
                    
                    if sample_imgs:
                        inputs = processor(text=feature_pool, images=sample_imgs, return_tensors="pt", padding=True)
                        with torch.no_grad():
                            outputs = model(**inputs)
                        
                        probs = outputs.logits_per_image.softmax(dim=1).mean(dim=0)
                        top_indices = probs.argsort(descending=True)[:4] # 가장 일치하는 특징 상위 4개
                        dynamic_features = [feature_pool[i] for i in top_indices]
                except:
                    pass
                
                if dynamic_features:
                    st.caption(f"🔍 **수집된 이미지 실제 특징 분석:** {', '.join(dynamic_features)}")
                
                category_dynamic_data[category] = {"features": dynamic_features}
                st.write("---")
                
            # CMF 종합 분석 섹션 추가
            st.header("🎨 종합 CMF (Color, Material, Finish) 가이드")
            st.write("컬러뿐만 아니라 수집된 레퍼런스 전체를 관통하는 **주요 재질(Material)**과 **표면 마감(Finish)** 트렌드 파악 결과입니다.")
            
            with st.spinner("이미지 픽셀 및 조형 딥러닝(CLIP & KMeans) CMF 정밀 분석 중..."):
                global_materials = []
                global_finishes = []
                
                # 1. Material & Finish 분석 (CLIP)
                try:
                    cmf_pool = {
                        "Material": [
                            "aluminum or metallic material", "natural wood material", "glass or transparent material",
                            "soft fabric or textile", "premium leather material", "industrial concrete or stone", 
                            "matte plastic material", "machined steel or titanium"
                        ],
                        "Finish": [
                            "matte, non-reflective finish", "glossy, highly reflective finish", "rough textured finish",
                            "perfectly smooth finish", "anodized or brushed finish", "weathered or vintage distress finish",
                            "translucent or frosted finish"
                        ]
                    }
                    
                    sample_imgs = []
                    for filepath in saved_files[:15]: # 속도를 위해 전체 중 15장 샘플링
                        try:
                            sample_imgs.append(Image.open(filepath).convert("RGB"))
                        except: pass
                    
                    if sample_imgs:
                        # Material 평가
                        m_inputs = processor(text=cmf_pool["Material"], images=sample_imgs, return_tensors="pt", padding=True)
                        with torch.no_grad(): m_outputs = model(**m_inputs)
                        m_probs = m_outputs.logits_per_image.softmax(dim=1).mean(dim=0)
                        top_m = m_probs.argsort(descending=True)[:2]
                        global_materials = [cmf_pool["Material"][i].split(" material")[0].title() for i in top_m]
                        
                        # Finish 평가
                        f_inputs = processor(text=cmf_pool["Finish"], images=sample_imgs, return_tensors="pt", padding=True)
                        with torch.no_grad(): f_outputs = model(**f_inputs)
                        f_probs = f_outputs.logits_per_image.softmax(dim=1).mean(dim=0)
                        top_f = f_probs.argsort(descending=True)[:2]
                        global_finishes = [cmf_pool["Finish"][i].split(" finish")[0].title() for i in top_f]
                except Exception as e:
                    pass
                    
                # 2. Color 분석 (KMeans)
                try:
                    all_pixels = []
                    for filepath in saved_files:
                        try:
                            img = Image.open(filepath).convert("RGB")
                            img = img.resize((50, 50)) # 빠른 처리를 위해 50x50 리사이즈
                            pixels_3d = np.array(img)
                            
                            # 컴팩트하고 빠른 배경 제거 (테두리 색상 추정법)
                            border_pixels = np.concatenate([
                                pixels_3d[0, :], pixels_3d[-1, :], 
                                pixels_3d[:, 0], pixels_3d[:, -1]
                            ])
                            bg_color = np.median(border_pixels, axis=0)
                            
                            flat_pixels = pixels_3d.reshape(-1, 3)
                            distances = np.linalg.norm(flat_pixels - bg_color, axis=1)
                            product_pixels = flat_pixels[distances > 40]
                            if len(product_pixels) < 100:
                                product_pixels = pixels_3d[12:38, 12:38].reshape(-1, 3)
                                
                            all_pixels.append(product_pixels)
                        except:
                            pass
                    
                    if all_pixels:
                        all_pixels = np.vstack(all_pixels)
                        if len(all_pixels) > 100000:
                            indices = np.random.choice(len(all_pixels), 100000, replace=False)
                            all_pixels = all_pixels[indices]
                            
                        num_colors = 8
                        kmeans = KMeans(n_clusters=num_colors, random_state=42, n_init='auto')
                        kmeans.fit(all_pixels)
                        colors = kmeans.cluster_centers_.astype(int)
                        labels, counts = np.unique(kmeans.labels_, return_counts=True)
                        colors = colors[np.argsort(counts)[::-1]]
                        
                        # CMF 텍스트 결과 먼저 표시
                        if global_materials or global_finishes:
                            m_text = ', '.join(global_materials) if global_materials else '분석 불가'
                            f_text = ', '.join(global_finishes) if global_finishes else '분석 불가'
                            st.success(f"**🛠️ 가장 지배적인 제질 (Material):** {m_text}  \n**✨ 제품 마감 트렌드 (Finish):** {f_text}")
                        
                        st.caption("🎨 **추출된 8대 대표 색상 (Color)**")
                        cols = st.columns(num_colors)
                        for i in range(num_colors):
                            hex_color = '#{:02x}{:02x}{:02x}'.format(colors[i][0], colors[i][1], colors[i][2]).upper()
                            with cols[i]:
                                st.markdown(f'''
                                    <div style="background-color: {hex_color}; height: 80px; border-radius: 12px; margin-bottom: 10px; border: 1px solid rgba(0,0,0,0.1); box-shadow: 0 4px 6px rgba(0,0,0,0.05);"></div>
                                ''', unsafe_allow_html=True)
                                st.code(hex_color, language="text")
                except Exception as e:
                    st.warning(f"색상 추출 중 오류가 발생했습니다: {e}")
                    
            st.write("---")
                
            # 프롬프트 생성 및 추천 섹션 추가
            st.header("✍️ AI 제품 디자인 프롬프트 추천")
            st.write("사전에 입력된 고정 단어가 아닌, 방금 전 핀터레스트에서 수집된 **실제 이미지들의 조형적 특징과 색상**을 AI가 딥러닝으로 분석하여, 제미나이가 원본 사진과 동일한 형태학적 질감을 띄도록 동적 프롬프트를 생성했습니다.")
            
            if len(active_categories) > 0:
                st.info("💡 단순히 외운 단어를 나열하는 것이 아니라, 수집된 이미지의 '실제 형태'와 '색상'을 기반으로 시안을 생성하도록 설계되었습니다.")
                for category in active_categories:
                    st.subheader(f"✨ {category} 렌더링 프롬프트")
                    
                    # 1. 해당 카테고리에 분류된 실제 이미지를 실시간 분석하여 대표 색상 3가지 추출
                    cat_hex_colors = []
                    try:
                        cat_pixels = []
                        # 속도를 위해 최대 10장의 이미지만 샘플링 분석
                        for filepath in classified_images[category][:10]: 
                            img = Image.open(filepath).convert("RGB").resize((30, 30))
                            pixels_3d = np.array(img)
                            border_pixels = np.concatenate([pixels_3d[0,:], pixels_3d[-1,:], pixels_3d[:,0], pixels_3d[:,-1]])
                            bg_color = np.median(border_pixels, axis=0)
                            flat = pixels_3d.reshape(-1, 3)
                            distances = np.linalg.norm(flat - bg_color, axis=1)
                            prod_pixels = flat[distances > 40]
                            if len(prod_pixels) < 50:
                                prod_pixels = pixels_3d[10:20, 10:20].reshape(-1, 3)
                            cat_pixels.append(prod_pixels)
                        
                        if cat_pixels:
                            cat_pixels = np.vstack(cat_pixels)
                            if len(cat_pixels) > 10000:
                                cat_pixels = cat_pixels[np.random.choice(len(cat_pixels), 10000, replace=False)]
                            km = KMeans(n_clusters=3, random_state=42, n_init='auto')
                            km.fit(cat_pixels)
                            for center in km.cluster_centers_.astype(int):
                                cat_hex_colors.append('#{:02x}{:02x}{:02x}'.format(*center).upper())
                    except:
                        pass
                    
                    # 2. 수집된 사진으로부터 동적 추출된 디자인 특징 구문 결합
                    features = category_dynamic_data.get(category, {}).get("features", [])
                    if features:
                        features_str = ", ".join(features)
                        desc_sentence = f"The product must visually translate these key physical elements extracted directly from the collected image data: {features_str}."
                    else:
                        desc_sentence = f"The design should perfectly capture the core physical essence of the {category} style."
                        
                    color_prompt = f"Strictly apply a color palette based on these extracted hex codes: {', '.join(cat_hex_colors)}." if cat_hex_colors else ""
                    
                    suggested_prompt = (
                        f"A photorealistic, highly detailed 3D actual product design of a '{keyword}', styled in a {category} aesthetic. "
                        f"The product must be isolated on a pure white background (background completely removed) to emphasize its form. "
                        f"{desc_sentence} "
                        f"{color_prompt} "
                        f"Render it exactly as an industrial design prototype, shot on a 50mm lens, 8k resolution, extremely professional studio product photography."
                    )
                    st.code(suggested_prompt, language="text")
                    
                    if category in category_dynamic_data:
                        category_dynamic_data[category]["hex_colors"] = cat_hex_colors
                        category_dynamic_data[category]["prompt"] = suggested_prompt
                        
            st.write("---")
            st.header("📥 PPTX 자동 완성 문서 다운로드")
            st.write("방금 앱에 나타난 무드보드 이미지, 조형 특징 추출 결과, 컬러 팔레트와 프롬프트까지 분석된 모든 자료를 클릭 단 한 번만으로 깔끔하게 정리된 파워포인트 문서(.pptx)로 저장할 수 있습니다.")
            
            try:
                pptx_stream = create_pptx_report(keyword, category_dynamic_data, classified_images)
                st.download_button(
                    label="📊 분석 리포트 다운로드 (.pptx)",
                    data=pptx_stream,
                    file_name=f"Design_Analysis_{keyword.replace(' ', '_')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"문서 생성 중 오류: {e}")

    except Exception as e:
        st.error(f"화면 출력 중 오류: {e}")